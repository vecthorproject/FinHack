import os
import io 
import pandas as pd 
import numpy as np 
import matplotlib.pyplot as plt 
import geopandas as gpd
import re 
import warnings 
from pptx import Presentation 
from pptx.util import Inches, Pt, Cm
import matplotlib.patheffects as pe

warnings.filterwarnings("ignore", category=UserWarning, module="matplotlib")
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

# 🟢 FIX 1: Formattazione sicura a prova di bomba per "n.d."
def format_euro(numero, decimali=2):
    if pd.isna(numero): return "n.d."
    try:
        num_float = float(numero)
        formato = f"{{:,.{decimali}f}}".format(num_float)
        return formato.replace(',', 'X').replace('.', ',').replace('X', '.')
    except (ValueError, TypeError):
        return "n.d."

# =================================================================
# ⚙️ FUNZIONI DI REPLACEMENT PER POWERPOINT
# =================================================================

def replace_text_in_paragraphs(paragraphs, context):
    """Elabora i paragrafi preservando i colori ereditati e pulendo i caratteri sporchi (es. \x0b)"""
    import re
    from pptx.dml.color import RGBColor
    
    color_map = {
        '[COLOR_A]': RGBColor(0, 176, 80),   # Verde
        '[COLOR_B]': RGBColor(237, 125, 49), # Arancione/Giallo
        '[COLOR_C]': RGBColor(255, 0, 0)     # Rosso
    }
    
    for paragraph in paragraphs:
        testo_intero = paragraph.text
        if "{{" in testo_intero and "}}" in testo_intero:
            modificato = False
            for chiave, valore in context.items():
                pattern = r"\{\{\s*" + re.escape(chiave) + r"\s*\}\}"
                if re.search(pattern, testo_intero):
                    testo_intero = re.sub(pattern, str(valore), testo_intero)
                    modificato = True
            
            if modificato:
                # 🧹 PULIZIA ESTREMA: Rimuove l'A capo morbido (\x0b) che crea scritte strane su PPTX
                testo_intero = testo_intero.replace('\x0b', ' ').replace('\r', '').replace('\\x0b', ' ')

                # CASO 1: Se ci sono tag [COLOR_] (I Rating A, B, C) usiamo la formattazione frammentata
                if any(c in testo_intero for c in color_map.keys()):
                    base_font_name = paragraph.runs[0].font.name if paragraph.runs else None
                    base_font_size = paragraph.runs[0].font.size if paragraph.runs else None
                    base_font_bold = paragraph.runs[0].font.bold if paragraph.runs else None
                    try:
                        base_font_color = paragraph.runs[0].font.color.rgb if paragraph.runs and paragraph.runs[0].font.color.type == 1 else None
                    except:
                        base_font_color = None
                        
                    parts = re.split(r'(\[COLOR_[ABC]\])', testo_intero)
                    paragraph.clear()
                    
                    for part in parts:
                        if not part: continue
                        run = paragraph.add_run()
                        is_color_tag = part in color_map
                        run.text = part.replace('[COLOR_A]', 'A').replace('[COLOR_B]', 'B').replace('[COLOR_C]', 'C')
                        
                        if base_font_name: run.font.name = base_font_name
                        if base_font_size: run.font.size = base_font_size
                        if base_font_bold is not None: run.font.bold = base_font_bold
                        
                        if is_color_tag:
                            run.font.color.rgb = color_map[part]
                            run.font.bold = True
                        elif base_font_color:
                            run.font.color.rgb = base_font_color
                            
                # CASO 2: Testo normale (Numeri ISTAT, Titoli, NACE). Metodo ultra-sicuro che non altera MAI i colori base!
                else:
                    if len(paragraph.runs) > 0:
                        paragraph.runs[0].text = testo_intero
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = "" # Svuota i frammenti rotti
                    else:
                        paragraph.text = testo_intero

def elabora_shape_per_testo(shape, context):
    """Cerca i testi {{}} nelle caselle, nelle tabelle e nei gruppi geometrici"""
    if shape.has_text_frame:
        replace_text_in_paragraphs(shape.text_frame.paragraphs, context)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_text_in_paragraphs(cell.text_frame.paragraphs, context)
    # Se è un gruppo, entra dentro il gruppo a cercare i testi
    if getattr(shape, "shape_type", None) == 6:
        for s in shape.shapes:
            elabora_shape_per_testo(s, context)

def get_shape_and_coords(shapes, placeholder):
    """Ricerca ricorsiva anti-errore per trovare la casella esatta anche se raggruppata o scritta male"""
    placeholder_clean = placeholder.replace(" ", "").lower()
    for shape in shapes:
        if shape.has_text_frame:
            # Rimuoviamo spazi, a capo e caratteri invisibili di PPTX
            testo = shape.text.replace(" ", "").replace("\n", "").replace("\r", "").replace("\x0b", "").lower()
            if placeholder_clean in testo:
                return shape, (shape.left, shape.top, shape.width, shape.height)
        # Se è un Gruppo, entra dentro e cerca!
        if getattr(shape, "shape_type", None) == 6:
            found_shape, coords = get_shape_and_coords(shape.shapes, placeholder)
            if found_shape:
                return found_shape, coords
    return None, None

def replace_placeholder_with_picture(slide, placeholder_text, image_stream):
    """Sostituisce il rettangolo fantasma con l'immagine, MANTENENDO IL LIVELLO Z ORIGINALE"""
    shape_da_eliminare, coords = get_shape_and_coords(slide.shapes, placeholder_text)
    
    if coords and shape_da_eliminare:
        x, y, cx, cy = coords
        image_stream.seek(0)
        
        # 1. Inserisce l'immagine (che PPTX mette automaticamente in primissimo piano)
        nuova_img = slide.shapes.add_picture(image_stream, x, y, cx, cy)
        
        # 2. TRUCCO MAGICO: Sposta il nodo XML della nuova immagine ESATTAMENTE PRIMA del placeholder
        # Questo garantisce che la mappa rimanga "dietro" ai tuoi boxini di testo!
        shape_da_eliminare.element.addprevious(nuova_img.element)
        
        # 3. Elimina fisicamente il rettangolo originale
        sp = shape_da_eliminare.element
        sp.getparent().remove(sp)

def get_trend_style(trend_word):
    """Icona e colore coerenti col verso del trend, usati nei box commento dei grafici andamento"""
    if trend_word in ('in crescita', 'in miglioramento'):
        return '▲', '16A34A'   # verde
    if trend_word in ('in contrazione', 'in peggioramento'):
        return '▼', 'DC2626'   # rosso
    return '●', '64748B'       # grigio (stabile)

def formatta_box_commento_grafico(slide, placeholder, icona_metrica, titolo, trend_word, testo):
    """Trasforma il box {{commento_grafico_*}} da testo piatto e centrato in una card leggibile:
    barra colorata in alto (coerente col trend) + icona della metrica + titolo per esteso
    (a capo centrato se lungo) + badge trend + testo di analisi."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    shape, coords = get_shape_and_coords(slide.shapes, placeholder)
    if not shape or not coords:
        return
    left, top, width, height = coords

    icona_trend, colore_hex = get_trend_style(trend_word)
    colore = RGBColor.from_string(colore_hex)

    # Barra accento in cima alla card, colorata in base al verso del trend
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(45720))
    barra.fill.solid()
    barra.fill.fore_color.rgb = colore
    barra.line.fill.background()
    barra.shadow.inherit = False

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Icona grande della metrica: riempie lo spazio vuoto della card e la rende identificabile a colpo d'occhio
    p_icona = tf.paragraphs[0]
    p_icona.alignment = PP_ALIGN.CENTER
    p_icona.space_after = Pt(10)
    r_icona = p_icona.add_run()
    r_icona.text = icona_metrica
    r_icona.font.size = Pt(44)

    p_titolo = tf.add_paragraph()
    p_titolo.alignment = PP_ALIGN.CENTER
    p_titolo.space_after = Pt(8)
    r_titolo = p_titolo.add_run()
    r_titolo.text = titolo.upper()
    r_titolo.font.bold = True
    r_titolo.font.size = Pt(18)
    r_titolo.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    p_badge = tf.add_paragraph()
    p_badge.alignment = PP_ALIGN.CENTER
    p_badge.space_after = Pt(14)
    r_badge = p_badge.add_run()
    r_badge.text = f"{icona_trend}  {trend_word.capitalize()}"
    r_badge.font.bold = True
    r_badge.font.size = Pt(16)
    r_badge.font.color.rgb = colore

    p_corpo = tf.add_paragraph()
    p_corpo.alignment = PP_ALIGN.CENTER
    r_corpo = p_corpo.add_run()
    r_corpo.text = testo
    r_corpo.font.size = Pt(14)
    r_corpo.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

# =================================================================
# 🤖 MOTORE NARRATIVO PER POWERPOINT
# =================================================================

def get_short_eco(rating):
    if rating == 'A': return "L'azienda si posiziona oltre il 2° terzile, superando la media dei competitor. Esprime una buona capacità di generare profitti operativi."
    elif rating == 'B': return "Performance in linea con il mercato. La marginalità è adeguata, con fisiologici spazi di ottimizzazione."
    return "Valori inferiori al 1° terzile. I costi comprimono la marginalità operativa rispetto agli standard di settore, richiedendo un intervento."

def get_short_patr(rating):
    if rating == 'A': return "Struttura solida e ben capitalizzata. Gli indici confermano un perfetto allineamento agli standard ottimali del mercato."
    elif rating == 'B': return "Struttura equilibrata e in linea con il mercato. Adeguata la copertura degli asset a lungo termine."
    return "Capitalizzazione limitata e forte dipendenza dal debito. Necessario un ribilanciamento urgente delle fonti di finanziamento."

def get_short_fin(rating):
    if rating == 'A': return "Gestione della liquidità ottimale. Generazione di cassa abbondante per coprire agevolmente gli impegni a breve."
    elif rating == 'B': return "Gestione della liquidità bilanciata e stabile. La società dimostra una serena e adeguata capacità di coprire gli impegni a breve termine."
    return "Criticità nella solvibilità a breve. Liquidità vincolata e potenziale tensione operativa sui pagamenti correnti."

def get_bullet_nazionale(rat_eco, rat_patr, rat_fin):
    b1 = "Si distingue per una buona redditività operativa rispetto al panel nazionale." if rat_eco == 'A' else ("Mantiene un posizionamento economico strutturalmente competitivo nel panel." if rat_eco == 'B' else "Mostra una debolezza nella redditività operativa rispetto al panel nazionale.")
    b2 = "Capitalizzazione solida e buona indipendenza finanziaria dai terzi." if rat_patr == 'A' else ("Struttura delle fonti equilibrata e in linea con le medie del settore." if rat_patr == 'B' else "Mostra una fragilità patrimoniale rispetto ai top performer del settore.")

    if rat_patr == 'C':
        b3 = "Margini di miglioramento focalizzati sul contenimento dell'indebitamento."
    elif rat_fin == 'A':
        b3 = "Buona capacità di generazione di cassa e liquidità immediata."
    elif rat_fin == 'B':
        b3 = "Solvibilità a breve termine adeguata e gestione della cassa bilanciata."
    else:
        b3 = "Tensioni di liquidità nel breve termine da monitorare attentamente."
    return b1, b2, b3

def get_bullet_regionale(rat_eco, rat_patr, rat_fin):
    b1 = "Si posiziona tra le realtà più solide del territorio sul piano industriale." if rat_eco == 'A' else ("Rappresenta una solida realtà industriale all'interno del proprio territorio." if rat_eco == 'B' else "Necessita di recuperare competitività industriale all'interno del mercato locale.")
    b2 = "Garantisce una buona flessibilità finanziaria rispetto ai competitor locali." if rat_fin == 'A' else ("Garantisce un equilibrio finanziario nel complesso adeguato." if rat_fin == 'B' else "Dinamiche di liquidità meno elastiche rispetto alla media delle imprese regionali.")
    b3 = "Presenta una struttura patrimoniale di riferimento a livello regionale." if rat_patr == 'A' else ("Allineamento fisiologico della struttura delle fonti rispetto ai competitor locali." if rat_patr == 'B' else "Evidenzia una minore capitalizzazione relativa rispetto ai competitor regionali più solidi.")
    return b1, b2, b3

def calcola_forza_debolezza(rating_eco, rating_patr, rating_fin):
    voti = {'Economico': rating_eco, 'Patrimoniale': rating_patr, 'Finanziario': rating_fin}
    ordinate = sorted(voti.items(), key=lambda x: x[1])
    migliore = ordinate[0]
    peggiore = ordinate[-1]

    # NOTA: il grado (A/B/C) del "migliore" e del "peggiore" va sempre verificato:
    # se tutte e tre le aree condividono lo stesso rating, il "migliore" relativo
    # non è necessariamente positivo (idem per il "peggiore"). Il testo deve
    # riflettere il rating effettivo, non solo la posizione relativa tra le tre aree.
    if migliore[1] == 'A':
        forza_titolo = f"Punto di forza in ambito {migliore[0]} (Classe {migliore[1]})"
        forza_testo = f"L'azienda presenta un posizionamento superiore alla mediana di settore nell'equilibrio {migliore[0].lower()}, a supporto della stabilità strategica."
    elif migliore[1] == 'B':
        forza_titolo = f"Area più solida: {migliore[0]} (Classe {migliore[1]})"
        forza_testo = f"Rispetto alle altre due aree, l'equilibrio {migliore[0].lower()} risulta il più allineato ai parametri mediani di settore."
    else:
        forza_titolo = f"Area relativamente meno critica: {migliore[0]} (Classe {migliore[1]})"
        forza_testo = f"Anche l'equilibrio {migliore[0].lower()}, pur essendo il meno critico tra i tre, si colloca al di sotto dei parametri mediani di settore: il quadro complessivo richiede attenzione su tutti i fronti."

    if peggiore[1] == 'C':
        att_titolo = f"Area di attenzione: {peggiore[0]} (Classe {peggiore[1]})"
        att_testo = f"Si rilevano le maggiori criticità nell'area di tipo {peggiore[0].lower()}. È opportuno concentrare le azioni di intervento su questo fronte."
    elif peggiore[1] == 'B':
        att_titolo = f"Area con margini di miglioramento: {peggiore[0]} (Classe {peggiore[1]})"
        att_testo = f"L'area di tipo {peggiore[0].lower()} presenta margini di ottimizzazione rispetto alle altre due, pur restando in linea con i parametri di settore."
    else:
        att_titolo = f"Area relativamente meno favorevole: {peggiore[0]} (Classe {peggiore[1]})"
        att_testo = f"Anche l'equilibrio {peggiore[0].lower()}, pur risultando il meno favorevole tra i tre, si mantiene superiore ai parametri mediani di settore."

    if peggiore[1] == 'C' and peggiore[0] == 'Patrimoniale':
        att_titolo = "Deleveraging strutturale da valutare (Classe C)"
        att_testo = "L'esposizione debitoria risulta superiore ai parametri di settore. È opportuno un rafforzamento del capitale proprio per contenere le tensioni di lungo termine."
    return forza_titolo, forza_testo, att_titolo, att_testo

# =================================================================
# 📝 COMMENTI STANDARDIZZATI PER I GRAFICI (Trend e Barre)
# =================================================================

def get_trend_diretto(val_21, val_24, inverso=False):
    try:
        v21, v24 = float(val_21), float(val_24)
        diff = v24 - v21
        # Variazione percentuale sul valore di partenza (2021), non differenza assoluta:
        # una soglia fissa (es. 0,5) e' tarata sui margini (scala 0-100) ma appiattisce a
        # "stabile" i ratio di liquidita'/struttura (scala 0-3), dove uno spostamento di
        # 0,2-0,3 e' gia' un miglioramento ben visibile nel grafico andamento accanto.
        base = abs(v21) if abs(v21) > 0.01 else abs(v24)
        var_pct = (diff / base * 100) if base > 0.01 else 0.0
        if inverso:
            if var_pct > 10: return "in peggioramento"
            elif var_pct < -10: return "in miglioramento"
        else:
            if var_pct > 10: return "in crescita"
            elif var_pct < -10: return "in contrazione"
    except (TypeError, ValueError, ZeroDivisionError):
        pass
    return "stabile"

def _raffronto_2021(trend, val_21, unita=''):
    """Ancora l'aggettivo di trend al dato 2021 che lo giustifica: senza il valore di
    partenza, un "stabile" o un "in crescita" restano aggettivi isolati e possono sembrare
    in contraddizione con il grafico andamento mostrato accanto, che i due anni li mette
    a confronto per esteso."""
    if not trend:
        return ""
    if val_21 is None or pd.isna(val_21):
        return f", {trend}"
    return f", {trend} (dal {format_euro(val_21)}{unita} del 2021)"

def get_commento_ebitda_trend(az_val, ita_val, trend='', val_21=None):
    vs = "superiore" if az_val >= ita_val else "inferiore"
    t = _raffronto_2021(trend, val_21, '%')
    return f"L'EBITDA Margin si attesta al {format_euro(az_val)}%{t}, risultando {vs} alla mediana settoriale ({format_euro(ita_val)}%)."

def get_commento_ebit_trend(az_val, ita_val, trend='', val_21=None):
    vs = "superiore" if az_val >= ita_val else "inferiore"
    t = _raffronto_2021(trend, val_21, '%')
    return f"L'EBIT Margin si attesta al {format_euro(az_val)}%{t}, risultando {vs} alla mediana settoriale ({format_euro(ita_val)}%)."

def get_commento_profit_trend(az_val, ita_val, trend='', val_21=None):
    vs = "superiore" if az_val >= ita_val else "inferiore"
    t = _raffronto_2021(trend, val_21, '%')
    return f"Il Profit Margin si attesta al {format_euro(az_val)}%{t}, risultando {vs} alla mediana settoriale ({format_euro(ita_val)}%)."

def get_commento_str1_trend(az_val, ita_val, trend='', val_21=None):
    soglia = "al di sopra della soglia di sicurezza" if az_val >= 1 else "al di sotto della soglia di sicurezza"
    t = _raffronto_2021(trend, val_21)
    return f"L'Indice primario di struttura si attesta a {format_euro(az_val)}{t}, risultando {soglia} (valore target ≥ 1)."

def get_commento_str2_trend(az_val, ita_val, trend='', val_21=None):
    soglia = "in equilibrio strutturale" if az_val >= 1 else "sotto la soglia di equilibrio"
    t = _raffronto_2021(trend, val_21)
    return f"L'Indice secondario di struttura si attesta a {format_euro(az_val)}{t}, risultando {soglia} (valore target ≥ 1)."

def get_commento_gearing_trend(az_val, ita_val, trend='', val_21=None):
    vs = "inferiore" if az_val <= ita_val else "superiore"
    dipendenza = "contenuta" if az_val <= ita_val else "elevata"
    t = _raffronto_2021(trend, val_21, '%')
    return f"Il Gearing si attesta al {format_euro(az_val)}%{t}, risultando {vs} alla mediana ({format_euro(ita_val)}%) con dipendenza debitoria {dipendenza}."

def get_commento_cr_trend(az_val, ita_val, trend='', val_21=None):
    soglia = "in equilibrio corrente" if az_val >= 1 else "in tensione corrente"
    t = _raffronto_2021(trend, val_21)
    return f"Il Current Ratio si attesta a {format_euro(az_val)}{t}, indicando una situazione {soglia} (valore target ≥ 1)."

def get_commento_qr_trend(az_val, ita_val, trend='', val_21=None):
    copertura = "senza necessità di smobilizzo delle scorte" if az_val >= 1 else "con parziale dipendenza dal magazzino"
    t = _raffronto_2021(trend, val_21)
    return f"Il Quick Ratio si attesta a {format_euro(az_val)}{t}, garantendo la copertura dei debiti a breve {copertura}."

def get_commento_rotazione_trend(az_val, ita_val, trend='', val_21=None):
    vs = "superiore" if az_val >= ita_val else "inferiore"
    t = _raffronto_2021(trend, val_21)
    return f"La Rotazione del Capitale si attesta a {format_euro(az_val)}{t}, risultando {vs} alla mediana settoriale ({format_euro(ita_val)})."

def get_commento_barre_eco(az_ebitda, az_ebit, az_prof, ita_ebitda, ita_ebit, ita_prof, reg_ebitda, reg_ebit, reg_prof):
    metriche = [
        ('EBITDA Margin', az_ebitda, ita_ebitda, reg_ebitda),
        ('EBIT Margin',   az_ebit,   ita_ebit,   reg_ebit),
        ('Profit Margin', az_prof,   ita_prof,   reg_prof),
    ]
    sotto_naz = [(n, az, ita) for n, az, ita, _ in metriche if az < ita]
    sopra_naz = [(n, az, ita) for n, az, ita, _ in metriche if az >= ita]
    sotto_reg = sum(1 for _, az, _, reg in metriche if az < reg)

    if not sotto_naz:
        frase1 = (f"Tutti gli indicatori superano la mediana nazionale: "
                  f"EBITDA Margin {format_euro(az_ebitda)}% (naz. {format_euro(ita_ebitda)}%), "
                  f"EBIT Margin {format_euro(az_ebit)}% (naz. {format_euro(ita_ebit)}%), "
                  f"Profit Margin {format_euro(az_prof)}% (naz. {format_euro(ita_prof)}%).")
    elif not sopra_naz:
        sotto_str = "; ".join(f"{n}: {format_euro(az)}% vs naz. {format_euro(ita)}%" for n, az, ita in sotto_naz)
        frase1 = f"Tutti gli indicatori si collocano al di sotto della mediana nazionale — {sotto_str}."
    else:
        sopra_str = ", ".join(f"{n}: {format_euro(az)}%" for n, az, _ in sopra_naz)
        sotto_str = ", ".join(f"{n}: {format_euro(az)}% (vs naz. {format_euro(ita)}%)" for n, az, ita in sotto_naz)
        frase1 = f"{sopra_str} supera la mediana nazionale; {sotto_str} si colloca al di sotto del benchmark."

    if sotto_reg == 0:
        frase2 = "Il confronto regionale conferma il posizionamento competitivo superiore alla media territoriale su tutte le metriche."
    elif sotto_reg == 3:
        frase2 = "Anche rispetto alla mediana regionale tutti gli indicatori risultano inferiori, accentuando le criticità operative a livello territoriale."
    else:
        frase2 = "Il confronto regionale evidenzia un posizionamento differenziato rispetto ai competitor del territorio."

    return f"{frase1} {frase2}"

def get_commento_barre_patr(az_str1, az_str2, az_gear, ita_str1, ita_str2, ita_gear, reg_str1, reg_str2, reg_gear):
    soglia_str1 = "≥ 1, in equilibrio" if az_str1 >= 1 else "< 1, sotto soglia"
    soglia_str2 = "≥ 1, in equilibrio" if az_str2 >= 1 else "< 1, sotto soglia"
    gear_vs_naz = "inferiore" if az_gear <= ita_gear else "superiore"
    gear_vs_reg = "inferiore" if az_gear <= reg_gear else "superiore"
    str1_vs_reg = "superiore" if az_str1 >= reg_str1 else "inferiore"
    str2_vs_reg = "superiore" if az_str2 >= reg_str2 else "inferiore"

    frase1 = (f"Ind. Struttura 1°: {format_euro(az_str1)} ({soglia_str1}); "
              f"Ind. Struttura 2°: {format_euro(az_str2)} ({soglia_str2}); "
              f"Gearing {format_euro(az_gear)}% — {gear_vs_naz} alla mediana nazionale ({format_euro(ita_gear)}%).")

    if az_str1 >= 1 and az_str2 >= 1 and az_gear <= ita_gear:
        frase2 = (f"Il confronto regionale conferma la solidità patrimoniale: entrambi gli indici di struttura risultano "
                  f"{str1_vs_reg} alla media territoriale e il Gearing è {gear_vs_reg} alla mediana regionale ({format_euro(reg_gear)}%).")
    elif az_str1 < 1 or az_str2 < 1:
        frase2 = (f"A livello regionale gli indici di struttura risultano {str1_vs_reg} (1°) e {str2_vs_reg} (2°) "
                  f"rispetto alla mediana territoriale; il Gearing è {gear_vs_reg} alla media regionale ({format_euro(reg_gear)}%).")
    else:
        frase2 = (f"La copertura degli asset risulta adeguata, ma il Gearing è {gear_vs_reg} alla mediana regionale "
                  f"({format_euro(reg_gear)}%), evidenziando un'esposizione debitoria da monitorare nel confronto territoriale.")

    return f"{frase1} {frase2}"

def get_commento_barre_fin(az_cr, az_qr, az_rot, ita_cr, ita_qr, ita_rot, reg_cr, reg_qr, reg_rot):
    cr_soglia = "≥ 1, equilibrio corrente" if az_cr >= 1 else "< 1, tensione corrente"
    qr_soglia = "≥ 1, liquidità adeguata"  if az_qr >= 1 else "< 1, dipendenza dal magazzino"
    rot_vs_naz = "superiore" if az_rot >= ita_rot else "inferiore"
    cr_vs_reg  = "superiore" if az_cr  >= reg_cr  else "inferiore"
    qr_vs_reg  = "superiore" if az_qr  >= reg_qr  else "inferiore"
    rot_vs_reg = "superiore" if az_rot >= reg_rot else "inferiore"

    frase1 = (f"Current Ratio {format_euro(az_cr)} ({cr_soglia}); "
              f"Quick Ratio {format_euro(az_qr)} ({qr_soglia}); "
              f"Rotazione Cap. Inv. {format_euro(az_rot)} — {rot_vs_naz} alla mediana naz. ({format_euro(ita_rot)}).")

    if az_cr >= 1 and az_qr >= 1 and az_rot >= ita_rot:
        frase2 = (f"Il confronto regionale conferma il posizionamento positivo: Current Ratio {cr_vs_reg} e "
                  f"Quick Ratio {qr_vs_reg} alla media territoriale, con rotazione del capitale {rot_vs_reg} "
                  f"rispetto alla mediana regionale ({format_euro(reg_rot)}).")
    elif not (az_cr >= 1) and not (az_qr >= 1):
        frase2 = (f"Il confronto regionale accentua le criticità di liquidità: Current Ratio {cr_vs_reg} e "
                  f"Quick Ratio {qr_vs_reg} rispetto alla media territoriale; la rotazione risulta {rot_vs_reg} "
                  f"alla mediana regionale ({format_euro(reg_rot)}).")
    else:
        frase2 = (f"A livello regionale, Current Ratio è {cr_vs_reg} e Quick Ratio è {qr_vs_reg} rispetto alla "
                  f"mediana territoriale; la rotazione del capitale è {rot_vs_reg} alla media regionale ({format_euro(reg_rot)}).")

    return f"{frase1} {frase2}"

# =================================================================
# 🚀 CORE GENERATION POWERPOINT
# =================================================================

def genera_presentazione_ppt(template_path, azienda_target, df_orbis, settore_nace, num_max_soc_orbis):
    df_raw = df_orbis.copy()
    
    # 🟢 FIX 2: PULIZIA INIZIALE ESTESA A TUTTI GLI ANNI E A TUTTE LE METRICHE
    base_numeriche = [
        'Totale valore della produzione migl EUR', 'Totale Attivo migl EUR', 'Numero dipendenti',
        'Margine di Profitto (*) %', 'Margine EBITDA (*) %', 'Margine EBIT (*) %',
        'Indice di Struttura 1° livello (*)', 'Indice di Struttura 2° livello (*)',
        'Gearing (*) %', 'Current Ratio (*)', 'Quick Ratio (*)', 'Indice di Rotazione del Capitale Investito (*)'
    ]
    
    anni_col = ['2021', '2022', '2023', '2024', '']
    for base in base_numeriche:
        for anno in anni_col:
            c = f"{base} {anno}".strip()
            if c in df_raw.columns:
                # Forza il rimpiazzo di testacci brutti di Excel in NaN e li converte
                df_raw[c] = pd.to_numeric(df_raw[c].astype(str).replace(['n.d.', 'n.a.', 'n.s.', 'N.D.', ' ', 'n.m.', 'nan'], np.nan), errors='coerce')
            
    col_ragione = [c for c in df_raw.columns if 'ragione' in str(c).lower()][0]
    col_nuts = [c for c in df_raw.columns if 'nuts2' in str(c).lower() or 'nuts 2' in str(c).lower()]
    col_regione = col_nuts[0] if col_nuts else None

    def punteggio_diretto(val, t1, t2):
        if pd.isna(val): return 3  # MODIFICATO: era return 1 — NaN → terzile peggiore
        return 1 if val >= t2 else (2 if val >= t1 else 3)  # MODIFICATO: era 3/.../1 — 1=primo terzile (best), 3=terzo terzile (worst)

    def punteggio_inverso(val, t1, t2):
        if pd.isna(val): return 3  # MODIFICATO: era return 1 — NaN → terzile peggiore
        return 1 if val <= t1 else (2 if val <= t2 else 3)  # MODIFICATO: era 3/.../1 — 1=primo terzile (best), 3=terzo terzile (worst)

    # MODIFICATO: due funzioni distinte perché assegna_lettera è usata in due contesti:
    # 1. assegna_lettera_area → pts_eco/fin/pat (somma 1/2/3 con 1=best, basso=buono)
    # 2. assegna_lettera      → pts_totali/Bench_Tot (valori_lettere A=3, alto=buono, invariata)
    def assegna_lettera_area(punti):  # NUOVO: soglia invertita, solo per somme di terzili area
        if pd.isna(punti): return 'C'
        return 'A' if punti <= 4 else ('B' if punti <= 7 else 'C')

    def assegna_lettera(punti):  # INVARIATA rispetto all'originale: usata per Bench_Tot
        if pd.isna(punti): return 'C'
        return 'A' if punti >= 8 else ('B' if punti >= 5 else 'C')

    # --- 1. IL MOTORE DI CALCOLO COMPLETO A 9 METRICHE (Copiatodal Word) ---
    c_prof = 'Margine di Profitto (*) % 2024'
    c_ebitda = 'Margine EBITDA (*) % 2024'
    c_ebit = 'Margine EBIT (*) % 2024'
    c_rot = 'Indice di Rotazione del Capitale Investito (*) 2024'
    c_quick = 'Quick Ratio (*) 2024'
    c_curr = 'Current Ratio (*) 2024'
    c_str1 = 'Indice di Struttura 1° livello (*) 2024'
    c_str2 = 'Indice di Struttura 2° livello (*) 2024'
    c_gear = 'Gearing (*) % 2024'

    # La Rotazione (c_rot) è tra le dirette: più alto è, più punti prende!
    metriche_dirette = [c_prof, c_ebitda, c_ebit, c_rot, c_quick, c_curr, c_str1, c_str2]
    
    # Il Gearing (c_gear) resta da solo tra le inverse: più è basso, meglio è!
    metriche_inverse = [c_gear]

    for m in metriche_dirette + metriche_inverse:
        if m in df_raw.columns:
            t1, t2 = df_raw[m].quantile(1/3), df_raw[m].quantile(2/3)
            funz = punteggio_inverso if m in metriche_inverse else punteggio_diretto
            df_raw[f'pts_{m}'] = df_raw[m].apply(lambda x: funz(x, t1, t2))
        else:
            df_raw[f'pts_{m}'] = 3  # Fallback: colonna assente → terzo terzile (MODIFICATO: era 1)

    # Somma Punti per Area
    df_raw['pts_eco'] = df_raw[f'pts_{c_prof}'] + df_raw[f'pts_{c_ebitda}'] + df_raw[f'pts_{c_ebit}']
    df_raw['pts_fin'] = df_raw[f'pts_{c_rot}'] + df_raw[f'pts_{c_quick}'] + df_raw[f'pts_{c_curr}']
    df_raw['pts_pat'] = df_raw[f'pts_{c_str1}'] + df_raw[f'pts_{c_str2}'] + df_raw[f'pts_{c_gear}']

    # Assegnazione Lettere Ufficiali — usa assegna_lettera_area (soglia invertita)
    df_raw['Rating Economico'] = df_raw['pts_eco'].apply(assegna_lettera_area)
    df_raw['Rating Finanziario'] = df_raw['pts_fin'].apply(assegna_lettera_area)
    df_raw['Rating Patrimoniale'] = df_raw['pts_pat'].apply(assegna_lettera_area)

    # Benchmark Totale
    valori_lettere = {'A': 3, 'B': 2, 'C': 1}
    df_raw['pts_totali'] = df_raw['Rating Economico'].map(valori_lettere) + \
                           df_raw['Rating Finanziario'].map(valori_lettere) + \
                           df_raw['Rating Patrimoniale'].map(valori_lettere)
    df_raw['Bench_Tot'] = df_raw['pts_totali'].apply(assegna_lettera)

    # --- 2. ESTRAZIONE DATI AZIENDA E TERRITORIO ---
    df_target = df_raw[df_raw[col_ragione].astype(str).str.lower().str.contains(azienda_target.lower().strip(), na=False)]
    if df_target.empty:
        raise ValueError("Azienda target non trovata nel campione analizzato.")
    riga = df_target.iloc[0]
    
    regione_target_pulita = riga.get(col_regione, 'N.D.').split(' - ')[-1] if isinstance(riga.get(col_regione, 'N.D.'), str) else 'N.D.'
    tot_imprese_regione = len(df_raw[df_raw[col_regione] == riga.get(col_regione)]) if col_regione else 0
    df_regione = df_raw[df_raw[col_regione] == riga.get(col_regione)] if col_regione else pd.DataFrame()

    # 👑 ELEZIONE DEL MARKET LEADER: prima rating combinato AAA, poi più grande per fatturato e attivo
    # MODIFICATO: era sort by pts_totali → ora filtra AAA first, poi ordina per dimensione
    _col_prod = 'Totale valore della produzione migl EUR 2024'
    _col_att  = 'Totale Attivo migl EUR 2024'
    _pool_aaa = df_raw[
        (df_raw['Rating Economico'] == 'A') &
        (df_raw['Rating Finanziario'] == 'A') &
        (df_raw['Rating Patrimoniale'] == 'A')
    ]
    _pool = _pool_aaa if not _pool_aaa.empty else df_raw  # fallback: tutto il campione se nessuna AAA
    _sort_cols  = [c for c in [_col_prod, _col_att] if c in _pool.columns]
    if _sort_cols:
        idx_leader = _pool.sort_values(by=_sort_cols, ascending=False).index[0]
    else:
        idx_leader = _pool.index[0]
        
    market_leader = str(df_raw.loc[idx_leader, col_ragione]) if pd.notna(idx_leader) and pd.notna(df_raw.loc[idx_leader, col_ragione]) else "N.D."

    # 🎯 ESTRAZIONE DELLE LETTERE PER L'AZIENDA TARGET (Esatte e coincidenti col Word)
    rat_eco = riga.get('Rating Economico', 'C')
    rat_pat = riga.get('Rating Patrimoniale', 'C')
    rat_fin = riga.get('Rating Finanziario', 'C')

    rat_eco_c = f"[COLOR_{rat_eco}]"
    rat_fin_c = f"[COLOR_{rat_fin}]"
    rat_pat_c = f"[COLOR_{rat_pat}]"
    
    # Ordine Istituzionale Corretto: ECO + PATR + FIN
    rat_tot_c = f"[COLOR_{rat_eco}][COLOR_{rat_pat}][COLOR_{rat_fin}]"

    forza_titolo, forza_testo, att_titolo, att_testo = calcola_forza_debolezza(rat_eco, rat_pat, rat_fin)
    naz1, naz2, naz3 = get_bullet_nazionale(rat_eco, rat_pat, rat_fin)
    reg1, reg2, reg3 = get_bullet_regionale(rat_eco, rat_pat, rat_fin)


    # =================================================================
    # 📊 ESTRAZIONE DATI ISTAT DA SOTTO-CARTELLA (XLSX) CON SOMMA MULTI-NACE
    # =================================================================

    # Nome della cartella definitiva
    NOME_CARTELLA_ISTAT = "fileistat" 

    def pulisci_numero_istat(val):
        """Pulisce i numeri sporchi di ISTAT (es. '..' per dati mancanti) e li converte"""
        if pd.isna(val) or str(val).strip() in ['', '..', 'n.d.', 'N.D.']:
            return 0.0
        try:
            v = str(val).replace(',', '.') 
            return float(v)
        except ValueError:
            return 0.0

    def estrai_dati_istat(nace_stringa, cartella):
        """Legge i file XLSX ISTAT cercando il NACE in modo totale, senza limiti di intestazione."""
        risultati = {
            'tot_imprese': 0, 'tot_val_prod_mln': 0.0, 'tot_dipendenti': 0,
            'no_num': 0, 'ne_num': 0, 'ce_num': 0, 'su_num': 0
        }
        
        codici_estratti = re.findall(r'\b\d{2}\.?\d{1,2}\b', str(nace_stringa))
        nace_puliti = [c.replace('.', '') for c in codici_estratti if len(c.replace('.', '')) >= 3]
        tags_nace = [f"[{n[:4]}]" for n in nace_puliti]
        
        if not tags_nace:
            return risultati
            
        file_soc = os.path.join(cartella, "TotSocIST.xlsx")
        file_dip = os.path.join(cartella, "TotDipIST.xlsx")
        file_val = os.path.join(cartella, "TotValEcoIST.xlsx")

        def trova_riga_sicura(df):
            # Analizza ogni riga unita come testo per trovare il tag NACE, indipendentemente dalla colonna
            for i, riga in df.iterrows():
                testo_riga = " ".join(riga.fillna("").astype(str))
                for tag in tags_nace:
                    if tag in testo_riga: return riga
                for n in nace_puliti:
                    if f"{n} " in testo_riga or f"{n[:2]}.{n[2:]}" in testo_riga: return riga
            return None

        # 1. LETTURA SOCIETA'
        if os.path.exists(file_soc):
            try:
                # Leggiamo tutto senza skiprows per evitare di saltare i dati
                df_soc = pd.read_excel(file_soc, header=None, dtype=str)
                riga = trova_riga_sicura(df_soc)
                if riga is not None:
                    # Assumiamo che i dati siano in colonna B(1), C(2), D(3), E(4), F(5), G(6)
                    risultati['tot_imprese'] += int(pulisci_numero_istat(riga.iloc[1]))
                    risultati['no_num'] += int(pulisci_numero_istat(riga.iloc[2]))
                    risultati['ne_num'] += int(pulisci_numero_istat(riga.iloc[3]))
                    risultati['ce_num'] += int(pulisci_numero_istat(riga.iloc[4]))
                    
                    sud = pulisci_numero_istat(riga.iloc[5])
                    isole = pulisci_numero_istat(riga.iloc[6])
                    risultati['su_num'] += int(sud + isole)
            except Exception as e:
                pass

        # 2. LETTURA DIPENDENTI
        if os.path.exists(file_dip):
            try:
                df_dip = pd.read_excel(file_dip, header=None, dtype=str)
                riga = trova_riga_sicura(df_dip)
                if riga is not None:
                    risultati['tot_dipendenti'] += int(pulisci_numero_istat(riga.iloc[1]))
            except Exception as e:
                pass

        # 3. LETTURA VALORE ECONOMICO
        if os.path.exists(file_val):
            try:
                df_val = pd.read_excel(file_val, header=None, dtype=str)
                riga = trova_riga_sicura(df_val)
                if riga is not None:
                    valore_migliaia = pulisci_numero_istat(riga.iloc[3])
                    risultati['tot_val_prod_mln'] += (valore_migliaia / 1000)
            except Exception as e:
                pass

        return risultati

    # Lanciamo l'estrazione
    dati_istat = estrai_dati_istat(settore_nace, NOME_CARTELLA_ISTAT)
    ha_dati_istat = dati_istat['tot_imprese'] > 0

    def calc_perc_istat(num, tot):
        if not ha_dati_istat or tot == 0: return "n.d."
        return f"{format_euro(num / tot * 100)}%"

    # Prepariamo le stringhe formattate con "n.d." se i file non ci sono o non trova nulla
    tot_is = dati_istat['tot_imprese']
    
    str_tot_imprese = f"{tot_is:,}".replace(',', '.') if ha_dati_istat else "n.d."
    str_tot_vprod = f"{dati_istat['tot_val_prod_mln']:,.0f}".replace(',', '.') if ha_dati_istat else "n.d."
    str_tot_dip = f"{dati_istat['tot_dipendenti']:,}".replace(',', '.') if ha_dati_istat else "n.d."
    
    str_no_num = f"{dati_istat['no_num']:,}".replace(',', '.') if ha_dati_istat else "n.d."
    str_ne_num = f"{dati_istat['ne_num']:,}".replace(',', '.') if ha_dati_istat else "n.d."
    str_ce_num = f"{dati_istat['ce_num']:,}".replace(',', '.') if ha_dati_istat else "n.d."
    str_su_num = f"{dati_istat['su_num']:,}".replace(',', '.') if ha_dati_istat else "n.d."

    istat_no_perc = calc_perc_istat(dati_istat['no_num'], tot_is)
    istat_ne_perc = calc_perc_istat(dati_istat['ne_num'], tot_is)
    istat_ce_perc = calc_perc_istat(dati_istat['ce_num'], tot_is)
    istat_su_perc = calc_perc_istat(dati_istat['su_num'], tot_is)

    # =================================================================
    # 📝 GENERAZIONE TESTO RILEVANZA TERRITORIALE
    # =================================================================
    aree_sort = sorted([
        ('Nord Ovest', dati_istat['no_num']),
        ('Nord Est', dati_istat['ne_num']),
        ('Centro', dati_istat['ce_num']),
        ('Sud e Isole', dati_istat['su_num'])
    ], key=lambda x: x[1], reverse=True)
    
    if ha_dati_istat and aree_sort[0][1] > 0:
        top1_area = aree_sort[0][0]
        top2_area = aree_sort[1][0] if aree_sort[1][1] > 0 else ""
        
        if top2_area:
            testo_rilevanza = f"Il panel presenta una copertura equilibrata sul territorio nazionale, con una maggiore concentrazione di imprese nel {top1_area} e nel {top2_area}, in termini di numero di società e valore economico."
        else:
            testo_rilevanza = f"Il panel presenta una forte concentrazione geografica, con la quasi totalità delle imprese localizzate nel {top1_area}."
    else:
        testo_rilevanza = "Dati di distribuzione territoriale non disponibili per questo settore."


    # =================================================================
    # 🏗️ COSTRUZIONE DIZIONARIO CONTESTO PPTX
    # =================================================================
    context = {
        'ragione_sociale': azienda_target,
        'codice_nace': str(settore_nace),
        'descr_settore': "", 
        'num_soc_valide': f"{len(df_raw):,}".replace(',', '.'),
        'tot_imprese_regione': f"{tot_imprese_regione:,}".replace(',', '.'),
        'regione_target': regione_target_pulita,
        # Valori diretti per i riquadri del Market Leader (Senza decimali)
        'tot_ricavi': format_euro(df_raw.loc[idx_leader].get('Totale valore della produzione migl EUR 2024', 0), 0) if pd.notna(idx_leader) else "n.d.",
        'tot_attivo': format_euro(df_raw.loc[idx_leader].get('Totale Attivo migl EUR 2024', 0), 0) if pd.notna(idx_leader) else "n.d.",
        'market_leader': market_leader,
        'rating_tot': rat_tot_c,
        'rating_eco': rat_eco_c,
        'rating_patr': rat_pat_c,
        'rating_fin': rat_fin_c,
        'descr_rating_eco_short': get_short_eco(rat_eco),
        'descr_rating_patr_short': get_short_patr(rat_pat),
        'descr_rating_fin_short': get_short_fin(rat_fin),
        'label_eco': "Forte" if rat_eco == 'A' else ("Adeguato" if rat_eco == 'B' else "Critico"),
        'label_patr': "Forte" if rat_pat == 'A' else ("Adeguato" if rat_pat == 'B' else "Critico"),
        'label_fin': "Forte" if rat_fin == 'A' else ("Adeguato" if rat_fin == 'B' else "Critico"),
        'naz_1': naz1, 'naz_2': naz2, 'naz_3': naz3,
        'reg_1': reg1, 'reg_2': reg2, 'reg_3': reg3,
        'punto_forza_titolo': forza_titolo,
        'punto_forza_testo': forza_testo,
        'area_attenzione_titolo': att_titolo,
        'area_attenzione_testo': att_testo,
        
        # 🟢 NUOVI TAG GEO ISTAT (Sicuri e protetti con "n.d.")
        'tot_imprese_istat': str_tot_imprese,
        'tot_vprod_is': str_tot_vprod,
        'tot_dip_istat': str_tot_dip,
        
        'no_num': str_no_num, 
        'no_perc': istat_no_perc,
        'ne_num': str_ne_num, 
        'ne_perc': istat_ne_perc,
        'ce_num': str_ce_num, 
        'ce_perc': istat_ce_perc,
        'su_num': str_su_num, 
        'su_perc': istat_su_perc,

        'testo_rilevanza': testo_rilevanza
    }

    # =================================================================
    # --- 3. GENERAZIONE GRAFICI MINI (TREND STORICI)
    # =================================================================
    def crea_grafico_mini(anni, valori_az, valori_set, valori_reg, titolo='', show_legend=False, is_percentage=False):
        fig, ax = plt.subplots(figsize=(3.6, 3.6))
        
        v_az = [float(v) if pd.notna(v) else 0.0 for v in valori_az]
        v_set = [float(v) if pd.notna(v) else 0.0 for v in valori_set]
        v_reg = [float(v) if pd.notna(v) else 0.0 for v in valori_reg]

        # Colori: Blu (Azienda), Arancio (Regione), Grigio scuro (Italia)
        col_az, col_reg, col_set = '#1D4ED8', '#F59E0B', '#1E293B'

        if titolo:
            ax.set_title(titolo, fontsize=10, fontweight='bold', color='#1E293B', pad=6) # 🟢 Titolo ingrandito

        # 🟢 LINEE E PUNTINI PIÙ SPESSI E VISIBILI
        ax.plot(anni, v_az, marker='o', linewidth=2.5, markersize=5, color=col_az, label='Azienda')
        ax.plot(anni, v_reg, marker='o', linewidth=2.0, markersize=4, color=col_reg, label=regione_target_pulita)
        ax.plot(anni, v_set, marker='o', linewidth=2.0, markersize=4, color=col_set, linestyle=':', label='Italia')

        ax.margins(y=0.5)

        suffix = "%" if is_percentage else ""
        
        # 🟢 ALGORITMO ANTI-SOVRAPPOSIZIONE
        for i in range(len(anni)):
            txt_az = format_euro(v_az[i]) + suffix
            txt_reg = format_euro(v_reg[i]) + suffix
            txt_set = format_euro(v_set[i]) + suffix

            # 1. Raggruppiamo i 3 punti dell'anno corrente
            punti = [
                {'val': v_az[i], 'txt': txt_az, 'color': col_az},
                {'val': v_reg[i], 'txt': txt_reg, 'color': col_reg},
                {'val': v_set[i], 'txt': txt_set, 'color': col_set}
            ]
            
            # 2. Li ordiniamo dal valore più ALTO al più BASSO
            punti.sort(key=lambda x: x['val'], reverse=True)
            
            # 3. Assegniamo offset Y scaglionati per distanziarli in verticale
            offsets = [(0, 11), (0, -14), (0, -27)]
            
            # 4. Stampiamo le etichette (più grandi e con bordo bianco)
            for j, p in enumerate(punti):
                ax.annotate(
                    p['txt'], 
                    (anni[i], p['val']), 
                    textcoords="offset points", 
                    xytext=offsets[j], 
                    ha='center', 
                    fontsize=9.0,          # 🟢 NUMERI INGRANDITI
                    fontweight='bold',  
                    color=p['color'],
                    path_effects=[pe.withStroke(linewidth=2.5, foreground='white')] # 🟢 BORDO BIANCO SALVAVITA
                )

        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.get_yaxis().set_visible(False)
        ax.spines['bottom'].set_color('#E2E8F0')

        # 🟢 LEGENDA PIÙ GRANDE
        ax.legend(loc='lower center', bbox_to_anchor=(0.5, -0.55), ncol=3, frameon=False, 
                fontsize=8.5, handletextpad=0.3, columnspacing=1.0)
                
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)

        mem_img = io.BytesIO()
        plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight')
        plt.close(fig)
        mem_img.seek(0)
        return mem_img
    
    anni = ['2021', '2022', '2023', '2024']
    def get_dati_grafico(nome_colonna):
        az = [riga.get(f'{nome_colonna} {a}', 0) for a in anni]
        set_m = [df_raw[f'{nome_colonna} {a}'].median() if f'{nome_colonna} {a}' in df_raw.columns else 0 for a in anni]
        reg_m = [df_regione[f'{nome_colonna} {a}'].median() if not df_regione.empty and f'{nome_colonna} {a}' in df_regione.columns else 0 for a in anni]
        return az, set_m, reg_m

    az, st, rg = get_dati_grafico('Margine di Profitto (*) %')
    img_eco_1 = crea_grafico_mini(anni, az, st, rg, titolo='Profit Margin', is_percentage=True)
    az, st, rg = get_dati_grafico('Margine EBIT (*) %')
    img_eco_2 = crea_grafico_mini(anni, az, st, rg, titolo='EBIT Margin', show_legend=True, is_percentage=True)
    az, st, rg = get_dati_grafico('Margine EBITDA (*) %')
    img_eco_3 = crea_grafico_mini(anni, az, st, rg, titolo='EBITDA Margin', is_percentage=True)

    az, st, rg = get_dati_grafico('Indice di Struttura 1° livello (*)')
    img_patr_1 = crea_grafico_mini(anni, az, st, rg, titolo='Ind. Strutt. 1° Liv.')
    az, st, rg = get_dati_grafico('Indice di Struttura 2° livello (*)')
    img_patr_2 = crea_grafico_mini(anni, az, st, rg, titolo='Ind. Strutt. 2° Liv.',show_legend=True)
    az, st, rg = get_dati_grafico('Gearing (*) %')
    img_patr_3 = crea_grafico_mini(anni, az, st, rg, titolo='Gearing', is_percentage=True)

    az, st, rg = get_dati_grafico('Current Ratio (*)')
    img_fin_1 = crea_grafico_mini(anni, az, st, rg, titolo='Current Ratio')
    az, st, rg = get_dati_grafico('Quick Ratio (*)')
    img_fin_2 = crea_grafico_mini(anni, az, st, rg, titolo='Quick Ratio', show_legend=True)
    az, st, rg = get_dati_grafico('Indice di Rotazione del Capitale Investito (*)')
    img_fin_3 = crea_grafico_mini(anni, az, st, rg, titolo='Rot. Cap. Inv.')
    # =================================================================
    # --- 4. GRAFICI A BARRE E TABELLE RANKING (2024)
    # =================================================================
    def crea_grafico_barre_confronto(metriche, val_ita, val_reg, val_az, nome_regione):
        fig, ax = plt.subplots(figsize=(7, 3.5))
        x = np.arange(len(metriche))
        width = 0.25
        
        v_ita = [float(v) if pd.notna(v) else 0.0 for v in val_ita]
        v_reg = [float(v) if pd.notna(v) else 0.0 for v in val_reg]
        v_az = [float(v) if pd.notna(v) else 0.0 for v in val_az]
        
        ax.bar(x - width, v_ita, width, label='Italia', color='#4F81BD', edgecolor='white')
        ax.bar(x, v_reg, width, label=nome_regione, color='#C0504D', edgecolor='white')
        ax.bar(x + width, v_az, width, label='Azienda', color='#9BBB59', edgecolor='white')
        
        ax.set_xticks(x)
        ax.set_xticklabels(metriche, fontsize=11, color='#595959')
        
        # 👇 MODIFICA: Spinta la legenda più in basso (da -0.2 a -0.35)
        ax.legend(loc='lower center', bbox_to_anchor=(0.5, -0.35), ncol=3, frameon=False, fontsize=10)
        
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['bottom'].set_color('#D9D9D9')
        ax.get_yaxis().set_visible(False)
        
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        mem_img = io.BytesIO()
        plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight')
        plt.close(fig)
        mem_img.seek(0)
        return mem_img

    def crea_tabella_confronto_img(dati_tabella, titolo_tabella):
        fig, ax = plt.subplots(figsize=(9, 1.8))
        ax.axis('tight')
        ax.axis('off')
        
        colonne = ['', 'Italia', 'Lazio', 'Azienda', 'Ranking Naz.', 'Ranking Reg.']
        table = ax.table(cellText=dati_tabella, colLabels=colonne, cellLoc='center', loc='center')
        table.auto_set_font_size(False)
        table.set_fontsize(10)
        table.scale(1.2, 1.8)
        table.auto_set_column_width([0])
        
        for (i, j), cell in table.get_celld().items():
            if i == 0:  
                cell.set_facecolor('#002060')
                cell.set_text_props(color='white', weight='bold')
            else:
                if j == 0: 
                    cell.set_facecolor('#002060')
                    cell.set_text_props(color='white', weight='bold', ha='left')

                    cell._text.set_position((0.1, 0.5)) 
                cell.set_edgecolor('#EAEAEA') 
                
        mem_img = io.BytesIO()
        plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight', transparent=True)
        plt.close(fig)
        mem_img.seek(0)
        return mem_img

    # 🟢 FIX 4: Comparazione Rank castata a Float per evitare crash su >, <
    def calc_rank_str(df_group, col_name, val_target, higher_is_better=True):
        if df_group.empty or col_name not in df_group.columns:
            return "n.d."
        
        # Usa il vero rank di pandas (esattamente come nel Word)
        asc_order = not higher_is_better
        rank_series = df_group[col_name].rank(ascending=asc_order, method='min')
        
        # Troviamo la riga esatta dell'azienda target
        riga_target = df_group[df_group[col_ragione].astype(str).str.lower().str.contains(azienda_target.lower().strip(), na=False)]
        
        if riga_target.empty or pd.isna(riga_target.iloc[0].name):
            return "n.d."
            
        idx_target = riga_target.iloc[0].name
        if pd.isna(rank_series.loc[idx_target]):
            return "n.d."
            
        pos_ottenuta = int(rank_series.loc[idx_target])
        tot_valide = df_group[col_name].notna().sum()
        return f"{format_euro(pos_ottenuta, 0)}/{format_euro(tot_valide, 0)}"

    # ECO
    az_ebitda, az_ebit, az_prof = riga.get('Margine EBITDA (*) % 2024', 0), riga.get('Margine EBIT (*) % 2024', 0), riga.get('Margine di Profitto (*) % 2024', 0)
    ita_ebitda = df_raw['Margine EBITDA (*) % 2024'].median() if 'Margine EBITDA (*) % 2024' in df_raw.columns else 0
    ita_ebit = df_raw['Margine EBIT (*) % 2024'].median() if 'Margine EBIT (*) % 2024' in df_raw.columns else 0
    ita_prof = df_raw['Margine di Profitto (*) % 2024'].median() if 'Margine di Profitto (*) % 2024' in df_raw.columns else 0
    reg_ebitda = df_regione['Margine EBITDA (*) % 2024'].median() if not df_regione.empty and 'Margine EBITDA (*) % 2024' in df_regione.columns else 0
    reg_ebit = df_regione['Margine EBIT (*) % 2024'].median() if not df_regione.empty and 'Margine EBIT (*) % 2024' in df_regione.columns else 0
    reg_prof = df_regione['Margine di Profitto (*) % 2024'].median() if not df_regione.empty and 'Margine di Profitto (*) % 2024' in df_regione.columns else 0

    img_barre_eco = crea_grafico_barre_confronto(['EBITDA Margin', 'EBIT Margin', 'Profit Margin'], [ita_ebitda, ita_ebit, ita_prof], [reg_ebitda, reg_ebit, reg_prof], [az_ebitda, az_ebit, az_prof], regione_target_pulita)
    rnk_naz_ebitda, rnk_reg_ebitda = calc_rank_str(df_raw, 'Margine EBITDA (*) % 2024', az_ebitda, True), calc_rank_str(df_regione, 'Margine EBITDA (*) % 2024', az_ebitda, True)
    rnk_naz_ebit, rnk_reg_ebit = calc_rank_str(df_raw, 'Margine EBIT (*) % 2024', az_ebit, True), calc_rank_str(df_regione, 'Margine EBIT (*) % 2024', az_ebit, True)
    rnk_naz_prof, rnk_reg_prof = calc_rank_str(df_raw, 'Margine di Profitto (*) % 2024', az_prof, True), calc_rank_str(df_regione, 'Margine di Profitto (*) % 2024', az_prof, True)

    img_tabella_eco = crea_tabella_confronto_img([
        ['EBITDA Margin %', format_euro(ita_ebitda), format_euro(reg_ebitda), format_euro(az_ebitda), rnk_naz_ebitda, rnk_reg_ebitda],
        ['EBIT Margin %', format_euro(ita_ebit), format_euro(reg_ebit), format_euro(az_ebit), rnk_naz_ebit, rnk_reg_ebit],
        ['Profit Margin %', format_euro(ita_prof), format_euro(reg_prof), format_euro(az_prof), rnk_naz_prof, rnk_reg_prof]
    ], "Equilibrio Economico - Anno 2024")

    # PATR
    az_str1, az_str2, az_gear = riga.get('Indice di Struttura 1° livello (*) 2024', 0), riga.get('Indice di Struttura 2° livello (*) 2024', 0), riga.get('Gearing (*) % 2024', 0)
    ita_str1 = df_raw['Indice di Struttura 1° livello (*) 2024'].median() if 'Indice di Struttura 1° livello (*) 2024' in df_raw.columns else 0
    ita_str2 = df_raw['Indice di Struttura 2° livello (*) 2024'].median() if 'Indice di Struttura 2° livello (*) 2024' in df_raw.columns else 0
    ita_gear = df_raw['Gearing (*) % 2024'].median() if 'Gearing (*) % 2024' in df_raw.columns else 0
    reg_str1 = df_regione['Indice di Struttura 1° livello (*) 2024'].median() if not df_regione.empty and 'Indice di Struttura 1° livello (*) 2024' in df_regione.columns else 0
    reg_str2 = df_regione['Indice di Struttura 2° livello (*) 2024'].median() if not df_regione.empty and 'Indice di Struttura 2° livello (*) 2024' in df_regione.columns else 0
    reg_gear = df_regione['Gearing (*) % 2024'].median() if not df_regione.empty and 'Gearing (*) % 2024' in df_regione.columns else 0

    # 👇 MODIFICA: Gearing rimosso dall'istogramma, mantenendo solo i due indici di struttura
    img_barre_patr = crea_grafico_barre_confronto(['Ind. Struttura 1°', 'Ind. Struttura 2°'], [ita_str1, ita_str2], [reg_str1, reg_str2], [az_str1, az_str2], regione_target_pulita)
    
    rnk_naz_str1, rnk_reg_str1 = calc_rank_str(df_raw, 'Indice di Struttura 1° livello (*) 2024', az_str1, True), calc_rank_str(df_regione, 'Indice di Struttura 1° livello (*) 2024', az_str1, True)
    rnk_naz_str2, rnk_reg_str2 = calc_rank_str(df_raw, 'Indice di Struttura 2° livello (*) 2024', az_str2, True), calc_rank_str(df_regione, 'Indice di Struttura 2° livello (*) 2024', az_str2, True)
    rnk_naz_gear, rnk_reg_gear = calc_rank_str(df_raw, 'Gearing (*) % 2024', az_gear, False), calc_rank_str(df_regione, 'Gearing (*) % 2024', az_gear, False)

    img_tabella_patr = crea_tabella_confronto_img([
        ['Ind. Struttura 1°', format_euro(ita_str1), format_euro(reg_str1), format_euro(az_str1), rnk_naz_str1, rnk_reg_str1],
        ['Ind. Struttura 2°', format_euro(ita_str2), format_euro(reg_str2), format_euro(az_str2), rnk_naz_str2, rnk_reg_str2],
        ['Gearing %', format_euro(ita_gear), format_euro(reg_gear), format_euro(az_gear), rnk_naz_gear, rnk_reg_gear]
    ], "Equilibrio Patrimoniale - Anno 2024")

    # FIN
    az_cr, az_qr, az_rot = riga.get('Current Ratio (*) 2024', 0), riga.get('Quick Ratio (*) 2024', 0), riga.get('Indice di Rotazione del Capitale Investito (*) 2024', 0)
    ita_cr = df_raw['Current Ratio (*) 2024'].median() if 'Current Ratio (*) 2024' in df_raw.columns else 0
    ita_qr = df_raw['Quick Ratio (*) 2024'].median() if 'Quick Ratio (*) 2024' in df_raw.columns else 0
    ita_rot = df_raw['Indice di Rotazione del Capitale Investito (*) 2024'].median() if 'Indice di Rotazione del Capitale Investito (*) 2024' in df_raw.columns else 0
    reg_cr = df_regione['Current Ratio (*) 2024'].median() if not df_regione.empty and 'Current Ratio (*) 2024' in df_regione.columns else 0
    reg_qr = df_regione['Quick Ratio (*) 2024'].median() if not df_regione.empty and 'Quick Ratio (*) 2024' in df_regione.columns else 0
    reg_rot = df_regione['Indice di Rotazione del Capitale Investito (*) 2024'].median() if not df_regione.empty and 'Indice di Rotazione del Capitale Investito (*) 2024' in df_regione.columns else 0

    img_barre_fin = crea_grafico_barre_confronto(['Current Ratio', 'Quick Ratio', 'Rotazione Cap.'], [ita_cr, ita_qr, ita_rot], [reg_cr, reg_qr, reg_rot], [az_cr, az_qr, az_rot], regione_target_pulita)
    rnk_naz_cr, rnk_reg_cr = calc_rank_str(df_raw, 'Current Ratio (*) 2024', az_cr, True), calc_rank_str(df_regione, 'Current Ratio (*) 2024', az_cr, True)
    rnk_naz_qr, rnk_reg_qr = calc_rank_str(df_raw, 'Quick Ratio (*) 2024', az_qr, True), calc_rank_str(df_regione, 'Quick Ratio (*) 2024', az_qr, True)
    rnk_naz_rot, rnk_reg_rot = calc_rank_str(df_raw, 'Indice di Rotazione del Capitale Investito (*) 2024', az_rot, True), calc_rank_str(df_regione, 'Indice di Rotazione del Capitale Investito (*) 2024', az_rot, True)

    img_tabella_fin = crea_tabella_confronto_img([
        ['Current Ratio', format_euro(ita_cr), format_euro(reg_cr), format_euro(az_cr), rnk_naz_cr, rnk_reg_cr],
        ['Quick Ratio', format_euro(ita_qr), format_euro(reg_qr), format_euro(az_qr), rnk_naz_qr, rnk_reg_qr],
        ['Rotazione Cap.', format_euro(ita_rot), format_euro(reg_rot), format_euro(az_rot), rnk_naz_rot, rnk_reg_rot]
    ], "Equilibrio Finanziario - Anno 2024")

    # =================================================================
    # 📝 COMMENTI STANDARDIZZATI PER I GRAFICI
    # =================================================================
    az_ebitda_21 = pd.to_numeric(riga.get('Margine EBITDA (*) % 2021', np.nan), errors='coerce')
    az_ebit_21   = pd.to_numeric(riga.get('Margine EBIT (*) % 2021', np.nan), errors='coerce')
    az_prof_21   = pd.to_numeric(riga.get('Margine di Profitto (*) % 2021', np.nan), errors='coerce')
    az_str1_21   = pd.to_numeric(riga.get('Indice di Struttura 1° livello (*) 2021', np.nan), errors='coerce')
    az_str2_21   = pd.to_numeric(riga.get('Indice di Struttura 2° livello (*) 2021', np.nan), errors='coerce')
    az_gear_21   = pd.to_numeric(riga.get('Gearing (*) % 2021', np.nan), errors='coerce')
    az_cr_21     = pd.to_numeric(riga.get('Current Ratio (*) 2021', np.nan), errors='coerce')
    az_qr_21     = pd.to_numeric(riga.get('Quick Ratio (*) 2021', np.nan), errors='coerce')
    az_rot_21    = pd.to_numeric(riga.get('Indice di Rotazione del Capitale Investito (*) 2021', np.nan), errors='coerce')

    # Trend calcolati una volta sola, riusati sia nel testo che nel badge colorato del box
    trend_eco_1, trend_eco_2, trend_eco_3 = get_trend_diretto(az_prof_21, az_prof), get_trend_diretto(az_ebit_21, az_ebit), get_trend_diretto(az_ebitda_21, az_ebitda)
    trend_patr_1, trend_patr_2, trend_patr_3 = get_trend_diretto(az_str1_21, az_str1), get_trend_diretto(az_str2_21, az_str2), get_trend_diretto(az_gear_21, az_gear, inverso=True)
    trend_fin_1, trend_fin_2, trend_fin_3 = get_trend_diretto(az_cr_21, az_cr), get_trend_diretto(az_qr_21, az_qr), get_trend_diretto(az_rot_21, az_rot)

    # ECO — img_eco_1=Profit, img_eco_2=EBIT, img_eco_3=EBITDA
    context['commento_barre_eco']     = get_commento_barre_eco(az_ebitda, az_ebit, az_prof, ita_ebitda, ita_ebit, ita_prof, reg_ebitda, reg_ebit, reg_prof)

    # PATR — img_patr_1=Strut1, img_patr_2=Strut2, img_patr_3=Gearing
    context['commento_barre_patr']     = get_commento_barre_patr(az_str1, az_str2, az_gear, ita_str1, ita_str2, ita_gear, reg_str1, reg_str2, reg_gear)

    # FIN — img_fin_1=CurrentRatio, img_fin_2=QuickRatio, img_fin_3=Rotazione
    context['commento_barre_fin']     = get_commento_barre_fin(az_cr, az_qr, az_rot, ita_cr, ita_qr, ita_rot, reg_cr, reg_qr, reg_rot)

    # 🎨 Box commento dei grafici andamento (icona + titolo per esteso + trend + testo): card, non {{}} piatti.
    # Formato tupla: (icona_metrica, titolo_per_esteso, trend, testo)
    dati_box_commenti = {
        'commento_grafico_eco_1': ('💰', 'Profit Margin', trend_eco_1, get_commento_profit_trend(az_prof, ita_prof, trend_eco_1, az_prof_21)),
        'commento_grafico_eco_2': ('📈', 'EBIT Margin', trend_eco_2, get_commento_ebit_trend(az_ebit, ita_ebit, trend_eco_2, az_ebit_21)),
        'commento_grafico_eco_3': ('📊', 'EBITDA Margin', trend_eco_3, get_commento_ebitda_trend(az_ebitda, ita_ebitda, trend_eco_3, az_ebitda_21)),
        'commento_grafico_patr_1': ('🏛️', 'Indice di Struttura 1° Livello', trend_patr_1, get_commento_str1_trend(az_str1, ita_str1, trend_patr_1, az_str1_21)),
        'commento_grafico_patr_2': ('🏗️', 'Indice di Struttura 2° Livello', trend_patr_2, get_commento_str2_trend(az_str2, ita_str2, trend_patr_2, az_str2_21)),
        'commento_grafico_patr_3': ('⚖️', 'Gearing', trend_patr_3, get_commento_gearing_trend(az_gear, ita_gear, trend_patr_3, az_gear_21)),
        'commento_grafico_fin_1': ('💧', 'Current Ratio', trend_fin_1, get_commento_cr_trend(az_cr, ita_cr, trend_fin_1, az_cr_21)),
        'commento_grafico_fin_2': ('⚡', 'Quick Ratio', trend_fin_2, get_commento_qr_trend(az_qr, ita_qr, trend_fin_2, az_qr_21)),
        'commento_grafico_fin_3': ('🔄', 'Indice di Rotazione del Capitale Investito', trend_fin_3, get_commento_rotazione_trend(az_rot, ita_rot, trend_fin_3, az_rot_21)),
    }

    # =================================================================
    # --- 5. TABELLA MARKET LEADER (SEMPLICE E PULITA) ---
    # =================================================================
    def crea_tabella_leader_colorata(dati, colonne, indici):
        fig, ax = plt.subplots(figsize=(6.8, 2))
        ax.axis('tight')
        ax.axis('off')
        
        table = ax.table(cellText=dati, rowLabels=indici, colLabels=colonne, cellLoc='center', loc='center')
        table.auto_set_font_size(False)
        table.set_fontsize(10)
        table.scale(1.0, 1.6)
        
        for (i, j), cell in table.get_celld().items():

            if i == 0:
                # Header row — font ridotto per intestazioni lunghe (es. "migl €")
                cell.set_facecolor('#002060')
                cell.set_text_props(color='white', weight='bold')
                cell.set_edgecolor('white')
                cell.set_fontsize(7.5)
            else:
                cell.set_facecolor('#F8FAFC')
                cell.set_edgecolor('#FFFFFF')
                if j == -1:
                    # Row labels column
                    cell.set_facecolor('#002060')
                    cell.set_text_props(color='white', weight='bold')
                    cell.set_edgecolor('white')
                elif j == 2:
                    cell.set_text_props(color='#002060', weight='bold', size=11)
                else:
                    cell.set_text_props(color='black')
        if (0, -1) in table.get_celld():
            table[0, -1].set_facecolor('#002060')
            table[0, -1].set_edgecolor('white')              
        mem_img = io.BytesIO()
        plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight', transparent=True)
        plt.close(fig)
        mem_img.seek(0)
        return mem_img

    # Estrazione Dati e Calcolo Rating del Leader
    if pd.notna(idx_leader):
        riga_leader = df_raw.loc[idx_leader]
        ml_ricavi = format_euro(riga_leader.get('Totale valore della produzione migl EUR 2024', 0), 0)
        ml_attivo = format_euro(riga_leader.get('Totale Attivo migl EUR 2024', 0), 0)
        # MODIFICATO: era == 3 (vecchia scala 3=best); ora 1=primo terzile=best, 3=terzo terzile=worst
        ml_eco = 'A' if riga_leader.get('pts_Margine EBITDA (*) % 2024', 3) == 1 else ('B' if riga_leader.get('pts_Margine EBITDA (*) % 2024', 3) == 2 else 'C')
        ml_pat = 'A' if riga_leader.get('pts_Gearing (*) % 2024', 3) == 1 else ('B' if riga_leader.get('pts_Gearing (*) % 2024', 3) == 2 else 'C')
        ml_fin = 'A' if riga_leader.get('pts_Current Ratio (*) 2024', 3) == 1 else ('B' if riga_leader.get('pts_Current Ratio (*) 2024', 3) == 2 else 'C')
        ml_bench = f"{ml_eco}{ml_pat}{ml_fin}" # Es: "AAA"
    else:
        ml_ricavi, ml_attivo, ml_bench = "n.d.", "n.d.", "N.D."

    # Valori Mediana Settore
    med_ricavi = format_euro(df_raw['Totale valore della produzione migl EUR 2024'].median(), 0)
    med_attivo = format_euro(df_raw['Totale Attivo migl EUR 2024'].median(), 0)
    med_bench = "BBB"

    # Valori Azienda Target
    az_ricavi = format_euro(riga.get('Totale valore della produzione migl EUR 2024', 0), 0)
    az_attivo = format_euro(riga.get('Totale Attivo migl EUR 2024', 0), 0)
    az_bench = f"{rat_eco}{rat_pat}{rat_fin}" # Es: "ABA"

    # Costruzione Matrice Tabella (Colonna Benchmark singola e pulita)
    colonne_tbl = ['Tot. Val. Prod. (migl €)', 'Tot. Attivo (migl €)', 'Benchmark']
    indici_tbl = [str(market_leader)[:35], str(azienda_target)[:35], 'Mediana Settore']
    
    dati_tbl = [
        [ml_ricavi, ml_attivo, ml_bench],
        [az_ricavi, az_attivo, az_bench],
        [med_ricavi, med_attivo, med_bench]
    ]

    img_tabella_leader = crea_tabella_leader_colorata(dati_tbl, colonne_tbl, indici_tbl)


    # =================================================================
    # 🗺️ GENERAZIONE MAPPA GEOPANDAS E CIAMBELLA CON LEGENDA (ANTI-CROP)
    # =================================================================
    def crea_grafico_ciambella(valori, etichette, percentuali, colori, totale_str):
        valori_puliti = []
        colori_puliti = []
        labels_puliti = []
        perc_pulite = []

        for v, c, e, p in zip(valori, colori, etichette, percentuali):
            if v > 0:
                valori_puliti.append(v)
                colori_puliti.append(c)
                labels_puliti.append(f"{e} ({p})")
                perc_pulite.append(p)

        fig, ax = plt.subplots(figsize=(7.0, 8.0))

        wedges = []
        if sum(valori_puliti) > 0:
            wedges, texts = ax.pie(
                valori_puliti,
                colors=colori_puliti,
                startangle=90,
                wedgeprops=dict(width=0.55, edgecolor='white', linewidth=2)
            )

            for i, (wedge, perc) in enumerate(zip(wedges, perc_pulite)):
                angle = (wedge.theta2 + wedge.theta1) / 2
                x = 0.72 * plt.np.cos(plt.np.radians(angle))
                y = 0.72 * plt.np.sin(plt.np.radians(angle))
                ax.text(x, y, perc, ha='center', va='center',
                        fontsize=13, fontweight='bold', color='white')

        ax.text(0, 0.12, totale_str, ha='center', va='center',
                fontsize=16, fontweight='bold', color='#1E293B')
        ax.text(0, -0.18, "Totale\nSocietà", ha='center', va='center',
                fontsize=11, color='#64748B', linespacing=1.4)

        lgd = None
        if wedges:
            lgd = ax.legend(wedges, labels_puliti,
                            loc="upper center",
                            bbox_to_anchor=(0.5, -0.08),
                            frameon=False,
                            fontsize=16,
                            handletextpad=0.6,
                            handlelength=1.2,
                            labelspacing=0.6)

        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        mem_img = io.BytesIO()

        if lgd:
            plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight', bbox_extra_artists=(lgd,))
        else:
            plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight')

        plt.close(fig)
        mem_img.seek(0)
        return mem_img

    def crea_mappa_italia_dinamica(n_no, n_ne, n_ce, n_su):
        url_geojson = "https://raw.githubusercontent.com/openpolis/geojson-italy/master/geojson/limits_IT_regions.geojson"
        try:
            mappa_italia = gpd.read_file(url_geojson)
        except:
            return None
            
        def mappa_reg_su_geo(nome_regione):
            r = str(nome_regione).lower()
            if any(x in r for x in ['piemonte', "valle d'aosta", "vallee", 'lombardia', 'liguria']): return 'Nord Ovest'
            if any(x in r for x in ['trentino', 'bolzano', 'bozen', 'veneto', 'friuli', 'emilia']): return 'Nord Est'
            if any(x in r for x in ['toscana', 'umbria', 'marche', 'lazio']): return 'Centro'
            if any(x in r for x in ['abruzzo', 'molise', 'campania', 'puglia', 'basilicata', 'calabria', 'sicilia', 'sardegna']): return 'Sud e Isole'
            return "Altro"
            
        mappa_italia['MacroArea'] = mappa_italia['reg_name'].apply(mappa_reg_su_geo)
        
        # Colora la zona SOLO se ISTAT dice che ci sono aziende, altrimenti grigio
        colori = {
            'Nord Ovest': '#0F172A' if n_no > 0 else '#E2E8F0', 
            'Nord Est': '#3B82F6' if n_ne > 0 else '#E2E8F0', 
            'Centro': '#65A30D' if n_ce > 0 else '#E2E8F0', 
            'Sud e Isole': '#F97316' if n_su > 0 else '#E2E8F0', 
            'Altro': '#E2E8F0'
        }
        mappa_italia['Colore'] = mappa_italia['MacroArea'].map(colori)
        
        fig, ax = plt.subplots(figsize=(6, 8))
        mappa_italia.plot(ax=ax, color=mappa_italia['Colore'], edgecolor='white', linewidth=1.2)
        ax.axis('off')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        mem_img = io.BytesIO()
        plt.savefig(mem_img, format='png', dpi=300, bbox_inches='tight')
        plt.close(fig)
        mem_img.seek(0)
        return mem_img

    # Passiamo 5 informazioni al grafico a ciambella!
    colori_geo = ['#0F172A', '#3B82F6', '#65A30D', '#F97316']
    etichette_geo = ['Nord Ovest', 'Nord Est', 'Centro', 'Sud e Isole']
    perc_geo = [istat_no_perc, istat_ne_perc, istat_ce_perc, istat_su_perc]
    valori_geo = [dati_istat['no_num'], dati_istat['ne_num'], dati_istat['ce_num'], dati_istat['su_num']]
    
    # Ecco la chiamata corretta con i 5 parametri in ordine:
    img_donut_geo = crea_grafico_ciambella(valori_geo, etichette_geo, perc_geo, colori_geo, str_tot_imprese)
    img_mappa_geo = crea_mappa_italia_dinamica(dati_istat['no_num'], dati_istat['ne_num'], dati_istat['ce_num'], dati_istat['su_num'])
    

    # =================================================================
    # 🖨️ ELABORAZIONE DEL FILE POWERPOINT
    # =================================================================
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            elabora_shape_per_testo(shape, context)
            
    # Sostituzione dei 9 Grafici Trend + 3 Grafici a Barre + 3 Tabelle Ranking + Tabella Leader
    for slide in prs.slides:
        replace_placeholder_with_picture(slide, "[[GRAFICO_ECO_1]]", img_eco_1)
        replace_placeholder_with_picture(slide, "[[GRAFICO_ECO_2]]", img_eco_2)
        replace_placeholder_with_picture(slide, "[[GRAFICO_ECO_3]]", img_eco_3)
        
        replace_placeholder_with_picture(slide, "[[GRAFICO_PATR_1]]", img_patr_1)
        replace_placeholder_with_picture(slide, "[[GRAFICO_PATR_2]]", img_patr_2)
        replace_placeholder_with_picture(slide, "[[GRAFICO_PATR_3]]", img_patr_3)
        
        replace_placeholder_with_picture(slide, "[[GRAFICO_FIN_1]]", img_fin_1)
        replace_placeholder_with_picture(slide, "[[GRAFICO_FIN_2]]", img_fin_2)
        replace_placeholder_with_picture(slide, "[[GRAFICO_FIN_3]]", img_fin_3)
        
        replace_placeholder_with_picture(slide, "[[GRAFICO_BARRE_ECO]]", img_barre_eco)
        replace_placeholder_with_picture(slide, "[[TABELLA_ECO_24]]", img_tabella_eco)
        
        replace_placeholder_with_picture(slide, "[[GRAFICO_BARRE_PATR]]", img_barre_patr)
        replace_placeholder_with_picture(slide, "[[TABELLA_PATR_24]]", img_tabella_patr)
        
        replace_placeholder_with_picture(slide, "[[GRAFICO_BARRE_FIN]]", img_barre_fin)
        replace_placeholder_with_picture(slide, "[[TABELLA_FIN_24]]", img_tabella_fin)
        
        replace_placeholder_with_picture(slide, "[[TABELLA_LEADER]]", img_tabella_leader)

        replace_placeholder_with_picture(slide, "[[GRAFICO_DONUT]]", img_donut_geo)
        if img_mappa_geo:
            replace_placeholder_with_picture(slide, "[[MAPPA_ITALIA]]", img_mappa_geo)

        for chiave, (icona_metrica, titolo, trend_word, testo) in dati_box_commenti.items():
            formatta_box_commento_grafico(slide, chiave, icona_metrica, titolo, trend_word, testo)

    output_ppt = io.BytesIO()
    prs.save(output_ppt)
    output_ppt.seek(0)
    
    return output_ppt


# =================================================================
# 📄 PDF EXPORT UTILITY
# =================================================================

def converti_pptx_in_pdf(pptx_bytes: "io.BytesIO | bytes") -> "io.BytesIO":
    """
    Converte un file PPTX (BytesIO o bytes) in PDF tramite LibreOffice.
    Restituisce un BytesIO contenente il PDF, oppure solleva RuntimeError se la
    conversione fallisce (LibreOffice non disponibile o errore interno).

    Utilizzo tipico:
        output_pptx = genera_presentazione_ppt(template_path, azienda, df, nace, n)
        output_pdf  = converti_pptx_in_pdf(output_pptx)
        with open("report.pdf", "wb") as f:
            f.write(output_pdf.read())
    """
    import subprocess
    import tempfile
    import shutil

    # --- Normalizza input ---
    if isinstance(pptx_bytes, (bytes, bytearray)):
        data = pptx_bytes
    else:
        pptx_bytes.seek(0)
        data = pptx_bytes.read()

    # --- Scrivi PPTX in cartella temp isolata ---
    tmp_dir = tempfile.mkdtemp(prefix="pptx2pdf_")
    try:
        pptx_path = os.path.join(tmp_dir, "report.pptx")
        with open(pptx_path, "wb") as f:
            f.write(data)

        # --- Lancia LibreOffice in headless ---
        soffice_candidates = ["soffice", "libreoffice", "/usr/bin/soffice", "/usr/bin/libreoffice"]
        soffice_bin = None
        for candidate in soffice_candidates:
            if shutil.which(candidate):
                soffice_bin = candidate
                break

        if soffice_bin is None:
            raise RuntimeError(
                "LibreOffice non trovato. Installarlo con: sudo apt-get install libreoffice\n"
                "Oppure usare converti_pptx_in_pdf_via_gotenberg() se disponibile nel tuo ambiente."
            )

        result = subprocess.run(
            [
                soffice_bin,
                "--headless",
                "--norestore",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice ha restituito codice {result.returncode}.\n"
                f"stderr: {result.stderr.strip()}"
            )

        pdf_path = pptx_path.replace(".pptx", ".pdf")
        if not os.path.exists(pdf_path):
            raise RuntimeError(
                f"LibreOffice ha completato senza errori ma il PDF non è stato generato.\n"
                f"stdout: {result.stdout.strip()}"
            )

        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        output = io.BytesIO(pdf_bytes)
        output.seek(0)
        return output

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def genera_presentazione_ppt_e_pdf(
    template_path: str,
    azienda_target: str,
    df_orbis,
    settore_nace: str,
    num_max_soc_orbis: int,
) -> "tuple[io.BytesIO, io.BytesIO]":
    """
    Wrapper conveniente: genera sia il PPTX che il PDF in una sola chiamata.

    Restituisce una tupla (pptx_bytes, pdf_bytes) — entrambi BytesIO seekati a 0.
    Se la conversione PDF fallisce, solleva RuntimeError con istruzioni.

    Esempio:
        pptx_io, pdf_io = genera_presentazione_ppt_e_pdf(template, azienda, df, nace, n)
        with open("report.pptx", "wb") as f: f.write(pptx_io.read())
        with open("report.pdf",  "wb") as f: f.write(pdf_io.read())
    """
    pptx_io = genera_presentazione_ppt(template_path, azienda_target, df_orbis, settore_nace, num_max_soc_orbis)
    pdf_io = converti_pptx_in_pdf(pptx_io)
    pptx_io.seek(0)  # reset dopo che converti_pptx_in_pdf l'ha letto
    return pptx_io, pdf_io
